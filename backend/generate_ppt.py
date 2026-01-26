"""
General Purpose PPT Generator
=============================
Converts content from PDF or text files into professional presentation slides.

Usage:
    1. Place your source file (PDF, TXT, or MD) in the 'input/' folder
    2. Update config.yaml with your settings
    3. Run: python generate_ppt.py
"""

import os
import re
import json
import yaml
import glob
from datetime import datetime
from pathlib import Path
from dotenv import load_dotenv

# Load environment variables
# Load environment variables
load_dotenv()

try:
    from prompts.title_prompts import get_title_slide_prompt
    from prompts.content_prompts import get_content_slide_prompt
    from prompts.ending_prompts import get_ending_slide_prompt
except ImportError:
    # Fallback/Development support if prompts folder isn't in path relatively
    import sys
    sys.path.append(str(Path(__file__).parent))
    from prompts.title_prompts import get_title_slide_prompt
    from prompts.content_prompts import get_content_slide_prompt
    from prompts.ending_prompts import get_ending_slide_prompt

# Get the directory where this script is located
SCRIPT_DIR = Path(__file__).parent.resolve()


def load_config():
    """Load configuration from config.yaml."""
    config_path = SCRIPT_DIR / "config.yaml"
    with open(config_path, "r") as f:
        return yaml.safe_load(f)


def get_instructions():
    """Reads the instructions.md file."""
    try:
        with open(SCRIPT_DIR / "instructions.md", "r") as f:
            return f.read()
    except FileNotFoundError:
        print("Warning: instructions.md not found, using default instructions.")
        return get_default_instructions()


def get_default_instructions():
    """Default slide design instructions."""
    return """
Design professional McKinsey-style presentation slides:
- Clean, minimalist design with clear hierarchy
- Use bullet points for lists
- Include data visualizations where appropriate
- Add "So What?" insight boxes for key takeaways
- Professional color scheme
"""


def extract_content_from_pdf(file_path, pages_to_process=-1):
    """Extract text content from a PDF file page by page."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        print("Error: PyMuPDF not installed. Run: pip install pymupdf")
        return []
    
    pages = []
    doc = fitz.open(file_path)
    
    total_pages = len(doc)
    if pages_to_process == -1 or pages_to_process > total_pages:
        pages_to_process = total_pages
    
    print(f"  Extracting {pages_to_process} pages from PDF...")
    
    for page_num in range(pages_to_process):
        page = doc.load_page(page_num)
        text = page.get_text("text")
        
        # Clean up the text
        text = text.strip()
        if text:
            pages.append({
                "page_number": page_num + 1,
                "content": text,
                "title": f"Page {page_num + 1}"
            })
    
    doc.close()
    return pages


def extract_content_from_text(file_path, pages_to_process=-1):
    """Extract content from a text or markdown file."""
    with open(file_path, "r", encoding="utf-8") as f:
        content = f.read()
    
    # Split by common page/section delimiters
    sections = []
    
    # Try markdown headers first (## or #)
    if re.search(r'^#{1,2}\s+', content, re.MULTILINE):
        parts = re.split(r'^(#{1,2}\s+.+)$', content, flags=re.MULTILINE)
        current_title = "Introduction"
        current_content = ""
        
        for i, part in enumerate(parts):
            if re.match(r'^#{1,2}\s+', part):
                if current_content.strip():
                    sections.append({
                        "title": current_title,
                        "content": current_content.strip()
                    })
                current_title = re.sub(r'^#{1,2}\s+', '', part).strip()
                current_content = ""
            else:
                current_content += part
        
        if current_content.strip():
            sections.append({
                "title": current_title,
                "content": current_content.strip()
            })
    
    # Try horizontal rules
    elif re.search(r'^-{3,}$|^={3,}$', content, re.MULTILINE):
        parts = re.split(r'\n-{3,}\n|\n={3,}\n', content)
        for i, part in enumerate(parts):
            if part.strip():
                sections.append({
                    "page_number": i + 1,
                    "title": f"Section {i + 1}",
                    "content": part.strip()
                })
    
    # Fall back to paragraphs (double newlines)
    else:
        paragraphs = re.split(r'\n\n+', content)
        chunk_size = 3
        for i in range(0, len(paragraphs), chunk_size):
            chunk = "\n\n".join(paragraphs[i:i + chunk_size])
            if chunk.strip():
                sections.append({
                    "page_number": (i // chunk_size) + 1,
                    "title": f"Section {(i // chunk_size) + 1}",
                    "content": chunk.strip()
                })
    
    # Apply page limit
    if pages_to_process != -1:
        sections = sections[:pages_to_process]
    
    # Add page numbers
    for i, section in enumerate(sections):
        section["page_number"] = i + 1
    
    return sections


def load_source_content(config):
    """Load content from the source file based on config."""
    input_config = config.get("input", {})
    file_name = input_config.get("file_name", "")
    pages_to_process = config.get("slides", {}).get("pages_to_process", -1)
    
    input_folder = SCRIPT_DIR / "input"
    file_path = input_folder / file_name
    
    if not file_path.exists():
        # Check if there's any file in input folder
        if input_folder.exists():
            files = [f for f in input_folder.glob("*.*") if not f.name.startswith('.')]
            if files:
                file_path = files[0]
                print(f"  Using found file: {file_path.name}")
            else:
                raise FileNotFoundError(f"No input file found in {input_folder}")
        else:
            raise FileNotFoundError(f"Input folder not found: {input_folder}")
    
    print(f"  Loading: {file_path.name}")
    
    # Determine content type from file extension
    ext = file_path.suffix.lower()
    if ext == ".pdf":
        return extract_content_from_pdf(file_path, pages_to_process)
    elif ext in [".txt", ".md", ".text", ".markdown"]:
        return extract_content_from_text(file_path, pages_to_process)
    else:
        return extract_content_from_text(file_path, pages_to_process)


def distribute_content_to_slides(pages, num_slides):
    """Distribute source pages across the specified number of slides."""
    if num_slides == -1 or num_slides >= len(pages):
        # One slide per page
        return [[page] for page in pages]
    
    if num_slides <= 0:
        num_slides = 1
    
    # Distribute pages across slides
    slides_content = [[] for _ in range(num_slides)]
    
    for i, page in enumerate(pages):
        slide_idx = i * num_slides // len(pages)
        if slide_idx >= num_slides:
            slide_idx = num_slides - 1
        slides_content[slide_idx].append(page)
    
    # Remove empty slides
    slides_content = [s for s in slides_content if s]
    
    return slides_content


def generate_title_slide_html(config, instructions):
    """Generate HTML for the title slide using AI."""
    title_config = config.get("title_slide", {})
    styling = config.get("content_styling", {})
    slide_config = config.get("slides", {})
    llm_config = config.get("llm", {})
    
    if not title_config.get("enabled", True):
        return None
    
    width = slide_config.get("width", 1280)
    height = slide_config.get("height", 720)
    
    title = title_config.get("title", "Presentation")
    subtitle = title_config.get("subtitle", "")
    author = title_config.get("author", "")
    
    date_text = title_config.get("date", "auto")
    if date_text == "auto":
        date_text = datetime.now().strftime("%B %Y")
    
    bg_color = title_config.get("background_color", "#004080")
    text_color = title_config.get("text_color", "#FFFFFF")
    font_family = title_config.get("font_family", "Inter")
    
    prompt = get_title_slide_prompt(
        title=title,
        subtitle=subtitle,
        author=author,
        date_text=date_text,
        width=width,
        height=height,
        bg_color=bg_color,
        text_color=text_color,
        styling=styling,
        font_family=font_family
    )

    provider = llm_config.get("provider", "claude")
    if provider == "claude":
        return _generate_with_claude(prompt, llm_config)
    else:
        return _generate_with_gemini(prompt, llm_config)


def generate_ending_slide_html(config, instructions):
    """Generate HTML for the ending slide using AI."""
    ending_config = config.get("ending_slide", {})
    styling = config.get("content_styling", {})
    slide_config = config.get("slides", {})
    llm_config = config.get("llm", {})
    
    if not ending_config.get("enabled", True):
        return None
    
    width = slide_config.get("width", 1280)
    height = slide_config.get("height", 720)
    bg_color = ending_config.get("background_color", "#004080")
    text_color = ending_config.get("text_color", "#FFFFFF")
    
    slide_type = ending_config.get("type", "thank_you")
    main_text = ending_config.get("main_text", "Thank You")
    secondary_text = ending_config.get("secondary_text", "")
    
    # Contact info
    contact_info = ""
    if slide_type == "contact":
        contact = ending_config.get("contact", {})
        contact_items = []
        if contact.get("email"):
            contact_items.append(f"Email: {contact['email']}")
        if contact.get("phone"):
            contact_items.append(f"Phone: {contact['phone']}")
        if contact.get("website"):
            contact_items.append(f"Website: {contact['website']}")
        if contact.get("address"):
            contact_items.append(contact['address'])
        contact_info = "\n".join(contact_items)
    
    if slide_type == "questions":
        main_text = main_text or "Questions?"
    
    prompt = get_ending_slide_prompt(
        slide_type=slide_type,
        main_text=main_text,
        secondary_text=secondary_text,
        contact_info=contact_info,
        width=width,
        height=height,
        bg_color=bg_color,
        text_color=text_color,
        styling=styling
    )

    provider = llm_config.get("provider", "claude")
    if provider == "claude":
        return _generate_with_claude(prompt, llm_config)
    else:
        return _generate_with_gemini(prompt, llm_config)


def generate_content_slide_html(slide_content, slide_number, total_slides, config, instructions):
    """Generate HTML for a content slide using LLM."""
    llm_config = config.get("llm", {})
    styling = config.get("content_styling", {})
    slide_config = config.get("slides", {})
    
    provider = llm_config.get("provider", "claude")
    
    width = slide_config.get("width", 1280)
    height = slide_config.get("height", 720)
    
    # Combine content from multiple pages if needed
    if isinstance(slide_content, list):
        combined_content = ""
        combined_title = ""
        for page in slide_content:
            if page.get("title"):
                if combined_title:
                    combined_title += " / "
                combined_title += page.get("title", "")
            combined_content += page.get("content", "") + "\n\n"
        content_str = combined_content.strip()
        page_title = combined_title or f"Slide {slide_number}"
    else:
        content_str = slide_content.get("content", "")
        page_title = slide_content.get("title", f"Slide {slide_number}")
    
    prompt = get_content_slide_prompt(
        page_title=page_title,
        content_str=content_str,
        slide_number=slide_number,
        total_slides=total_slides,
        instructions=instructions,
        styling=styling,
        width=width,
        height=height
    )

    if provider == "claude":
        return _generate_with_claude(prompt, llm_config)
    else:
        return _generate_with_gemini(prompt, llm_config)


def _generate_with_claude(prompt, llm_config):
    """Generate HTML using Claude API."""
    from anthropic import Anthropic
    
    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        raise ValueError("ANTHROPIC_API_KEY not found in environment")
    
    client = Anthropic(api_key=api_key)
    claude_config = llm_config.get("claude", {})
    
    message = client.messages.create(
        model=claude_config.get("model", "claude-sonnet-4-5-20250929"),
        max_tokens=claude_config.get("max_tokens", 10000),
        messages=[{"role": "user", "content": prompt}]
    )
    
    content_text = message.content[0].text.strip()
    
    # Clean up response
    if "```html" in content_text:
        content_text = content_text.split("```html")[1].split("```")[0].strip()
    elif "```" in content_text:
        content_text = content_text.split("```")[1].split("```")[0].strip()
    
    return content_text


def _generate_with_gemini(prompt, llm_config):
    """Generate HTML using Gemini API."""
    import google.generativeai as genai
    
    api_key = os.getenv("GOOGLE_API_KEY")
    if not api_key:
        raise ValueError("GOOGLE_API_KEY not found in environment")
    
    genai.configure(api_key=api_key)
    gemini_config = llm_config.get("gemini", {})
    
    model = genai.GenerativeModel(gemini_config.get("model", "gemini-3-pro-preview"))
    response = model.generate_content(prompt, request_options={"timeout": 1000})
    
    content_text = response.text.strip()
    
    # Clean up response
    if "```html" in content_text:
        content_text = content_text.split("```html")[1].split("```")[0].strip()
    elif "```" in content_text:
        content_text = content_text.split("```")[1].split("```")[0].strip()
    
    return content_text


def save_html_slide(html_content, slide_number, output_folder):
    """Save HTML content to a file."""
    file_path = output_folder / f"slide_{slide_number}.html"
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(html_content)
    return file_path


def convert_to_pdf(output_folder, config):
    """Convert all HTML slides to a single PDF using WeasyPrint (stable alternative to Playwright)."""
    from weasyprint import HTML, CSS
    from PyPDF2 import PdfMerger
    import tempfile
    
    slide_config = config.get("slides", {})
    output_config = config.get("output", {})
    
    width = slide_config.get("width", 1280)
    height = slide_config.get("height", 720)
    
    # Find all HTML files
    html_files = sorted(
        output_folder.glob("slide_*.html"),
        key=lambda x: int(re.search(r'slide_(\d+)', x.name).group(1))
    )
    
    if not html_files:
        print("  No HTML slides found to convert")
        return None
    
    print(f"\n  Converting {len(html_files)} slides to PDF using WeasyPrint...")
    
    # Output file name
    date_str = datetime.now().strftime("%Y-%m-%d")
    output_name = output_config.get("file_name", "Presentation")
    output_file = output_folder / f"{output_name}_{date_str}.pdf"
    
    temp_pdfs = []
    
    # CSS for proper sizing
    css_string = f"""
    @page {{
        size: {width}px {height}px;
        margin: 0;
    }}
    body {{
        margin: 0;
        padding: 0;
        width: {width}px;
        height: {height}px;
    }}
    """
    
    for i, html_file in enumerate(html_files, 1):
        print(f"    [{i}/{len(html_files)}] Converting: {html_file.name}")
        
        try:
            # Read HTML content
            with open(html_file, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            # Wrap with proper HTML structure if needed
            if '<!DOCTYPE' not in html_content.upper():
                html_content = f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{ 
            font-family: Arial, sans-serif;
            width: {width}px;
            height: {height}px;
            overflow: hidden;
        }}
    </style>
</head>
<body>
{html_content}
</body>
</html>'''
            
            # Create temporary PDF
            temp_pdf = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
            temp_pdfs.append(temp_pdf.name)
            temp_pdf.close()
            
            # Convert HTML to PDF
            html_doc = HTML(string=html_content)
            html_doc.write_pdf(temp_pdf.name, stylesheets=[CSS(string=css_string)])
            
        except Exception as e:
            print(f"    ✗ Failed to convert {html_file.name}: {str(e)}")
            raise
    
    # Merge all PDFs
    print("  Merging PDFs...")
    merger = PdfMerger()
    for temp_pdf in temp_pdfs:
        merger.append(temp_pdf)
    
    merger.write(str(output_file))
    merger.close()
    
    # Cleanup temp files
    for temp_pdf in temp_pdfs:
        try:
            os.unlink(temp_pdf)
        except:
            pass
    
    print(f"  ✓ PDF saved: {output_file}")
    return output_file


def convert_to_pptx(output_folder, config):
    """Convert all HTML slides to a PPTX file."""
    from bs4 import BeautifulSoup
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    
    slide_config = config.get("slides", {})
    output_config = config.get("output", {})
    
    width = slide_config.get("width", 1280)
    height = slide_config.get("height", 720)
    
    SCALE_FACTOR = Inches(13.333) / 1280
    
    html_files = sorted(
        output_folder.glob("slide_*.html"),
        key=lambda x: int(re.search(r'slide_(\d+)', x.name).group(1))
    )
    
    if not html_files:
        print("  No HTML slides found to convert")
        return None
    
    print(f"\n  Converting {len(html_files)} slides to PPTX...")
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for i, html_file in enumerate(html_files, 1):
        print(f"    [{i}/{len(html_files)}] Converting: {html_file.name}")
        
        with open(html_file, 'r') as f:
            html_content = f.read()
        
        soup = BeautifulSoup(html_content, 'html.parser')
        slide_div = soup.find('div', id='slide')
        
        if not slide_div:
            print(f"      Warning: No slide container found in {html_file.name}")
            continue
        
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        style = slide_div.get('style', '')
        bg_match = re.search(r'background-color:\s*(#[0-9a-fA-F]+)', style)
        if bg_match:
            hex_color = bg_match.group(1).lstrip('#')
            rgb = tuple(int(hex_color[j:j+2], 16) for j in (0, 2, 4))
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(*rgb)
        
        for element in slide_div.find_all(recursive=False):
            style = element.get('style', '')
            styles = {}
            for s in style.split(';'):
                if ':' in s:
                    k, v = s.split(':', 1)
                    styles[k.strip()] = v.strip()
            
            def parse_val(val):
                if not val:
                    return 0
                match = re.match(r"([\d\.]+)", val)
                return float(match.group(1)) if match else 0
            
            left = parse_val(styles.get('left', '0')) * SCALE_FACTOR
            top = parse_val(styles.get('top', '0')) * SCALE_FACTOR
            width_px = parse_val(styles.get('width', '100')) * SCALE_FACTOR
            height_px = parse_val(styles.get('height', '50')) * SCALE_FACTOR
            
            shape = slide.shapes.add_textbox(left, top, width_px, height_px)
            tf = shape.text_frame
            tf.word_wrap = True
            
            text = element.get_text(strip=True)
            if text:
                tf.paragraphs[0].text = text
                
                font_size = parse_val(styles.get('font-size', '11'))
                if font_size:
                    tf.paragraphs[0].font.size = Pt(font_size * 0.75)
                
                color = styles.get('color', '#000000')
                if color.startswith('#'):
                    hex_c = color.lstrip('#')
                    if len(hex_c) == 3:
                        hex_c = ''.join([c*2 for c in hex_c])
                    rgb = tuple(int(hex_c[j:j+2], 16) for j in (0, 2, 4))
                    tf.paragraphs[0].font.color.rgb = RGBColor(*rgb)
                
                if 'bold' in styles.get('font-weight', ''):
                    tf.paragraphs[0].font.bold = True
            
            bg_color = styles.get('background-color')
            if bg_color and bg_color.startswith('#'):
                hex_c = bg_color.lstrip('#')
                if len(hex_c) == 3:
                    hex_c = ''.join([c*2 for c in hex_c])
                rgb = tuple(int(hex_c[j:j+2], 16) for j in (0, 2, 4))
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*rgb)
    
    date_str = datetime.now().strftime("%Y-%m-%d")
    output_name = output_config.get("file_name", "Presentation")
    output_file = output_folder / f"{output_name}_{date_str}.pptx"
    
    prs.save(str(output_file))
    print(f"  ✓ PPTX saved: {output_file}")
    return output_file


def main():
    """Main entry point."""
    print("\n" + "=" * 60)
    print("  General Purpose PPT Generator")
    print("=" * 60)
    
    # Load configuration
    print("\n[1/5] Loading configuration...")
    config = load_config()
    
    # Create folders
    input_folder = SCRIPT_DIR / "input"
    output_folder = SCRIPT_DIR / "output"
    input_folder.mkdir(exist_ok=True)
    output_folder.mkdir(exist_ok=True)
    
    # Clean old HTML files from output
    for old_file in output_folder.glob("slide_*.html"):
        old_file.unlink()
    
    # Load instructions
    instructions = get_instructions()
    
    # Load source content
    print("\n[2/5] Loading source content...")
    try:
        pages = load_source_content(config)
        print(f"  ✓ Loaded {len(pages)} pages/sections from source")
    except FileNotFoundError as e:
        print(f"  ✗ Error: {e}")
        print(f"\n  Please place your source file in: {input_folder}")
        return
    
    if not pages:
        print("  ✗ No content found in source file")
        return
    
    # Distribute content across slides
    num_slides = config.get("slides", {}).get("number_of_slides", -1)
    slides_content = distribute_content_to_slides(pages, num_slides)
    print(f"  ✓ Distributing content across {len(slides_content)} content slides")
    
    # Generate slides
    print("\n[3/5] Generating slides with AI...")
    slide_number = 1
    all_html_slides = []
    
    processing_config = config.get("processing", {})
    max_retries = processing_config.get("max_retries", 2)
    
    # Title slide (AI generated)
    print(f"  Generating slide {slide_number}: Title slide...")
    for attempt in range(max_retries + 1):
        try:
            title_html = generate_title_slide_html(config, instructions)
            if title_html:
                save_html_slide(title_html, slide_number, output_folder)
                all_html_slides.append(title_html)
                print(f"  ✓ Slide {slide_number}: Title slide")
                slide_number += 1
            break
        except Exception as e:
            if attempt < max_retries:
                print(f"    Retry {attempt + 1}...")
            else:
                print(f"  ✗ Failed title slide: {e}")
    
    # Content slides
    total_content_slides = len(slides_content)
    for i, slide_content in enumerate(slides_content):
        print(f"  Generating slide {slide_number}: Content {i+1}/{total_content_slides}...")
        
        for attempt in range(max_retries + 1):
            try:
                html = generate_content_slide_html(
                    slide_content, 
                    slide_number, 
                    total_content_slides + 2,  # +2 for title and ending
                    config, 
                    instructions
                )
                save_html_slide(html, slide_number, output_folder)
                all_html_slides.append(html)
                print(f"  ✓ Slide {slide_number}: Content slide {i+1}")
                slide_number += 1
                break
            except Exception as e:
                if attempt < max_retries:
                    print(f"    Retry {attempt + 1}...")
                else:
                    print(f"  ✗ Failed slide {slide_number}: {e}")
    
    # Ending slide (AI generated)
    print(f"  Generating slide {slide_number}: Ending slide...")
    for attempt in range(max_retries + 1):
        try:
            ending_html = generate_ending_slide_html(config, instructions)
            if ending_html:
                save_html_slide(ending_html, slide_number, output_folder)
                all_html_slides.append(ending_html)
                print(f"  ✓ Slide {slide_number}: Ending slide")
            break
        except Exception as e:
            if attempt < max_retries:
                print(f"    Retry {attempt + 1}...")
            else:
                print(f"  ✗ Failed ending slide: {e}")
    
    # Convert to output format
    print("\n[4/5] Converting to output format...")
    output_format = config.get("output", {}).get("format", "pdf")
    
    if output_format == "pdf":
        output_file = convert_to_pdf(output_folder, config)
    else:
        output_file = convert_to_pptx(output_folder, config)
    
    # Cleanup
    print("\n[5/5] Cleanup...")
    if processing_config.get("cleanup_html", False):
        for html_file in output_folder.glob("slide_*.html"):
            html_file.unlink()
        print("  ✓ Removed intermediate HTML files")
    else:
        print("  ✓ HTML files preserved in output folder")
    
    # Summary
    print("\n" + "=" * 60)
    print("  Generation Complete!")
    print("=" * 60)
    print(f"  Total slides: {len(all_html_slides)}")
    if output_file:
        print(f"  Output file: {output_file}")
    print("=" * 60 + "\n")


if __name__ == "__main__":
    main()
